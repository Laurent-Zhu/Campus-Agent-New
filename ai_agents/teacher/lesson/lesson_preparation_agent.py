import os
from docx import Document
import PyPDF2
from pptx import Presentation
import openai
import json

from .agent_tools import (
    parse_document,
    generate_lesson_plan,
    generate_training_plan,
    generate_schedule,
    generate_ppt_outline,
)

client = openai.OpenAI(
    api_key=os.environ.get("DEEPSEEK_API_KEY"),
    base_url="https://api.deepseek.com/v1"
)

def parse_docx(file_path):
    """解析 docx 文件，返回全部文本内容"""
    doc = Document(file_path)
    content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return content

def parse_pdf(file_path):
    """解析 pdf 文件，返回全部文本内容"""
    content = ""
    with open(file_path, "rb") as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page in pdf_reader.pages:
            content += page.extract_text() or ""
    return content

def lesson_preparation_agent(file_path):
    # 1. 判断文件类型并解析
    if file_path.endswith('.docx'):
        content = parse_docx(file_path)
    elif file_path.endswith('.pdf'):
        content = parse_pdf(file_path)
    else:
        return {"error": "仅支持 docx 和 pdf 文件"}

    # 2. 构造 prompt
    prompt = (
        "你是一名智能教案助手。请根据以下课程材料内容，依次生成：\n"
        "1. 结构化课件草稿（包括课件标题、模块、要点等，JSON格式）\n"
        "2. 实训计划（包括实训目标、任务列表，JSON格式）\n"
        "3. 时间安排表（包括总课时、详细安排，JSON格式）\n"
        "4. 课件PPT大纲（包括每个幻灯片的标题和内容，JSON格式）\n"
        "请严格只输出JSON，不要加任何注释、解释或markdown代码块。\n"
        "{\n"
        "  \"structured_draft\": {...},\n"
        "  \"training_plan\": {...},\n"
        "  \"schedule\": {...}\n"
        "  \"ppt_outline\": {...}\n"
        "}\n"
        "课程材料内容如下：\n"
        f"{content[:3000]}"  # 控制长度，防止超长
    )
    
    # 3. 调用大模型
    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[
            {"role": "user", "content": prompt}
        ]
    )

    # 4. 返回大模型输出
    # return response.choices[0].message.content
    result_str = response.choices[0].message.content
    result_str = result_str.replace("```json", "").replace("```", "").strip()

    print("大模型原始输出：", result_str)
    try:
        result = json.loads(result_str)
    except Exception:
        return {"error": "大模型输出解析失败", "raw": result_str}
    return result

def generate_ppt_from_outline(ppt_outline: dict, outline_path: str):
    prs = Presentation()
    # Add a title slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = ppt_outline.get('title', '课件ppt')

    # Add slides for each section
    for slide_info in ppt_outline.get('slides', []):
        slide_layout = prs.slide_layouts[1]  # Content slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_info.get('title', '')
        slide.shapes.placeholders[1].text = slide_info.get('content', '')

    
    prs.save(outline_path)